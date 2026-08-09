from __future__ import annotations

import csv
import base64
import hashlib
import io
import json
import mimetypes
import os
import re
import sqlite3
import threading
import time
import uuid
import zipfile
import xml.etree.ElementTree as ET
from html import escape
from pathlib import Path
from typing import Any

from db_compat import connect_app_db

try:
    from docx import Document
except Exception:  # pragma: no cover - optional dependency
    Document = None  # type: ignore[assignment]

try:
    from openpyxl import Workbook, load_workbook
except Exception:  # pragma: no cover - optional dependency
    Workbook = None  # type: ignore[assignment]
    load_workbook = None  # type: ignore[assignment]

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - optional dependency
    Presentation = None  # type: ignore[assignment]

try:
    from pypdf import PdfReader
except Exception:  # pragma: no cover - optional dependency
    PdfReader = None  # type: ignore[assignment]

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
except Exception:  # pragma: no cover - optional dependency
    A4 = None  # type: ignore[assignment]
    ParagraphStyle = None  # type: ignore[assignment]
    getSampleStyleSheet = None  # type: ignore[assignment]
    mm = None  # type: ignore[assignment]
    pdfmetrics = None  # type: ignore[assignment]
    TTFont = None  # type: ignore[assignment]
    Paragraph = None  # type: ignore[assignment]
    SimpleDocTemplate = None  # type: ignore[assignment]
    Spacer = None  # type: ignore[assignment]

SUPPORTED_FORMATS = {"txt", "md", "json", "csv", "pdf", "docx", "xlsx", "pptx"}
OFFICE_FORMATS = {"docx", "xlsx", "pptx"}
TEXT_FORMATS = {"txt", "md", "json", "csv"}
MIME_BY_FORMAT = {
    "txt": "text/plain; charset=utf-8",
    "md": "text/markdown; charset=utf-8",
    "json": "application/json",
    "csv": "text/csv; charset=utf-8",
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
EXT_BY_FORMAT = {fmt: f".{fmt}" for fmt in SUPPORTED_FORMATS}
CONTROL_RE = re.compile(r"[\x00-\x1f\x7f]")


class ArtifactError(ValueError):
    pass


def now_ts() -> int:
    return int(time.time())


def clean_display_name(name: str, fallback: str = "file") -> str:
    name = CONTROL_RE.sub("", str(name or "")).strip().replace("\\", "/").split("/")[-1]
    name = name.strip(" .")[:180]
    return name or fallback


def normalize_format(name: str, requested: str = "") -> str:
    requested = str(requested or "").strip().lower().lstrip(".")
    if requested in SUPPORTED_FORMATS:
        return requested
    suffix = Path(clean_display_name(name)).suffix.lower().lstrip(".")
    if suffix in SUPPORTED_FORMATS:
        return suffix
    raise ArtifactError("unsupported file format")


def _xlsx_col_name(index: int) -> str:
    name = ""
    while index > 0:
        index, rem = divmod(index - 1, 26)
        name = chr(ord("A") + rem) + name
    return name or "A"


def _xlsx_cell_xml(cell_ref: str, value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, bool):
        return f'<c r="{cell_ref}" t="b"><v>{"1" if value else "0"}</v></c>'
    if isinstance(value, (int, float)) and not isinstance(value, bool):
        return f'<c r="{cell_ref}"><v>{value}</v></c>'
    text = escape(str(value))
    return f'<c r="{cell_ref}" t="inlineStr"><is><t xml:space="preserve">{text}</t></is></c>'


def _build_xlsx_bytes(content: Any) -> bytes:
    rows: list[list[Any]] = []
    if isinstance(content, dict) and isinstance(content.get("rows"), list):
        headers = content.get("headers")
        if isinstance(headers, list):
            rows.append(list(headers))
        for row in content["rows"]:
            rows.append(list(row) if isinstance(row, list) else [row])
    elif isinstance(content, list):
        for row in content:
            rows.append(list(row) if isinstance(row, list) else [row])
    else:
        for line in _plain_text_content(content).splitlines() or [""]:
            rows.append([line])
    sheet_rows = []
    for r_index, row in enumerate(rows, start=1):
        cells = []
        for c_index, value in enumerate(row, start=1):
            cell_ref = f"{_xlsx_col_name(c_index)}{r_index}"
            cell_xml = _xlsx_cell_xml(cell_ref, value)
            if cell_xml:
                cells.append(cell_xml)
        if cells:
            sheet_rows.append(f'<row r="{r_index}">{"".join(cells)}</row>')
        else:
            sheet_rows.append(f'<row r="{r_index}"/>')
    sheet_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        f"<sheetData>{''.join(sheet_rows) or '<row r=\"1\"/>'}</sheetData>"
        "</worksheet>"
    )
    workbook_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        '<sheets><sheet name="Data" sheetId="1" r:id="rId1"/></sheets>'
        "</workbook>"
    )
    rels_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/>'
        '<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles" Target="styles.xml"/>'
        '</Relationships>'
    )
    root_rels_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>'
        '</Relationships>'
    )
    content_types_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>'
        '<Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>'
        '<Override PartName="/xl/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.styles+xml"/>'
        '</Types>'
    )
    styles_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<styleSheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
        '<fonts count="1"><font><sz val="11"/><name val="Calibri"/></font></fonts>'
        '<fills count="1"><fill><patternFill patternType="none"/></fill></fills>'
        '<borders count="1"><border><left/><right/><top/><bottom/><diagonal/></border></borders>'
        '<cellStyleXfs count="1"><xf numFmtId="0" fontId="0" fillId="0" borderId="0"/></cellStyleXfs>'
        '<cellXfs count="1"><xf numFmtId="0" fontId="0" fillId="0" borderId="0" xfId="0"/></cellXfs>'
        '</styleSheet>'
    )
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", content_types_xml)
        zf.writestr("_rels/.rels", root_rels_xml)
        zf.writestr("xl/workbook.xml", workbook_xml)
        zf.writestr("xl/_rels/workbook.xml.rels", rels_xml)
        zf.writestr("xl/worksheets/sheet1.xml", sheet_xml)
        zf.writestr("xl/styles.xml", styles_xml)
    return buffer.getvalue()


def _extract_xlsx_text(path: Path, *, max_chars: int = 120_000) -> tuple[str, dict[str, Any]]:
    metadata: dict[str, Any] = {"sheets": ["Data"]}
    chunks: list[str] = []
    with zipfile.ZipFile(path) as zf:
        shared_strings: list[str] = []
        if "xl/sharedStrings.xml" in zf.namelist():
            try:
                root = ET.fromstring(zf.read("xl/sharedStrings.xml"))
                ns = {"main": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
                for si in root.findall("main:si", ns):
                    text_parts = [node.text or "" for node in si.findall(".//main:t", ns)]
                    shared_strings.append("".join(text_parts))
            except Exception:
                shared_strings = []
        sheet_name = "xl/worksheets/sheet1.xml"
        if sheet_name not in zf.namelist():
            raise ArtifactError("invalid workbook structure")
        root = ET.fromstring(zf.read(sheet_name))
        ns = {"main": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
        row_count = 0
        for row in root.findall(".//main:sheetData/main:row", ns):
            values: list[str] = []
            has_value = False
            for cell in row.findall("main:c", ns):
                ref = str(cell.attrib.get("r") or "")
                if not ref:
                    continue
                ctype = str(cell.attrib.get("t") or "")
                value = ""
                if ctype == "inlineStr":
                    value = "".join(node.text or "" for node in cell.findall(".//main:t", ns))
                elif ctype == "s":
                    idx_text = "".join(node.text or "0" for node in cell.findall("main:v", ns))
                    try:
                        idx = int(idx_text or "0")
                        value = shared_strings[idx] if 0 <= idx < len(shared_strings) else ""
                    except ValueError:
                        value = ""
                elif ctype == "b":
                    value = "TRUE" if "".join(node.text or "" for node in cell.findall("main:v", ns)) == "1" else "FALSE"
                else:
                    value = "".join(node.text or "" for node in cell.findall("main:v", ns))
                values.append(value)
                has_value = has_value or bool(value)
            if has_value:
                row_count += 1
                chunks.append(" | ".join(values))
            if len("\n".join(chunks)) >= max_chars:
                break
        metadata["sheet_rows"] = {"Data": row_count}
    text = "\n".join(chunks)
    if len(text) > max_chars:
        metadata["text_truncated"] = True
        text = text[:max_chars]
    metadata["extracted_chars"] = len(text)
    return text, metadata


def _pdf_escape_literal(text: str) -> str:
    return text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def _sanitize_pdf_visible_text(text: str) -> str:
    visible = re.sub(r"[\x00-\x1f\x7f]+", " ", str(text or "")).strip()
    visible = visible.encode("ascii", "ignore").decode("ascii", "ignore")
    return visible or "Personal Agent Rus report"


def _build_pdf_bytes(content: Any) -> bytes:
    raw_text = _plain_text_content(content)
    visible = _sanitize_pdf_visible_text(raw_text.splitlines()[0] if raw_text.strip() else raw_text)
    payload = base64.b64encode(raw_text.encode("utf-8")).decode("ascii")
    stream_text = f"BT /F1 12 Tf 72 770 Td ({_pdf_escape_literal(visible)}) Tj ET"
    objects = [
        "<< /Type /Catalog /Pages 2 0 R >>",
        "<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        "<< /Type /Page /Parent 2 0 R /MediaBox [0 0 595 842] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
        f"<< /Length {len(stream_text.encode('ascii'))} >>\nstream\n{stream_text}\nendstream",
        "<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    parts = [b"%PDF-1.4\n", f"%%PAR_TEXT:{payload}\n".encode("ascii")]
    offsets = [0]
    offset = len(parts[0]) + len(parts[1])
    for index, obj in enumerate(objects, start=1):
        chunk = f"{index} 0 obj\n{obj}\nendobj\n".encode("ascii")
        parts.append(chunk)
        offsets.append(offset)
        offset += len(chunk)
    xref_start = offset
    xref_lines = ["xref\n0 6\n0000000000 65535 f \n"]
    for obj_offset in offsets[1:]:
        xref_lines.append(f"{obj_offset:010d} 00000 n \n")
    trailer = f"trailer\n<< /Size 6 /Root 1 0 R >>\nstartxref\n{xref_start}\n%%EOF\n"
    parts.append("".join(xref_lines).encode("ascii"))
    parts.append(trailer.encode("ascii"))
    return b"".join(parts)


def _extract_pdf_text(path: Path, *, max_chars: int = 120_000) -> tuple[str, dict[str, Any]]:
    data = path.read_bytes()
    marker = re.search(rb"%%PAR_TEXT:([A-Za-z0-9+/=]+)", data)
    if marker:
        try:
            text = base64.b64decode(marker.group(1)).decode("utf-8", errors="replace")
        except Exception:
            text = ""
        metadata = {"pages": 1, "fallback": True, "extracted_chars": len(text)}
        if len(text) > max_chars:
            metadata["text_truncated"] = True
            text = text[:max_chars]
        return text, metadata
    return "", {"pages": 0, "fallback": True, "extracted_chars": 0}


def _build_fallback_office_bytes(content: Any, *, marker: str) -> bytes:
    raw_text = _plain_text_content(content)
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("content.txt", raw_text.encode("utf-8"))
        zf.writestr("metadata.json", json.dumps({"marker": marker, "content_length": len(raw_text)}, ensure_ascii=False).encode("utf-8"))
    return buffer.getvalue()


def _extract_fallback_office_text(path: Path, *, marker: str, max_chars: int = 120_000) -> tuple[str, dict[str, Any]]:
    with zipfile.ZipFile(path) as zf:
        data = b""
        if "content.txt" in zf.namelist():
            data = zf.read("content.txt")
        elif "metadata.json" in zf.namelist():
            data = zf.read("metadata.json")
        text = data.decode("utf-8", errors="replace")
    metadata = {"fallback": True, "marker": marker, "extracted_chars": 0}
    if len(text) > max_chars:
        metadata["text_truncated"] = True
        text = text[:max_chars]
    metadata["extracted_chars"] = len(text)
    return text, metadata


def _validate_zip_bytes(data: bytes, *, max_entries: int = 3000, max_uncompressed: int = 80 * 1024 * 1024, max_ratio: int = 200) -> None:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            infos = zf.infolist()
            if len(infos) > max_entries:
                raise ArtifactError("archive contains too many entries")
            total = 0
            for info in infos:
                normalized = info.filename.replace("\\", "/")
                parts = [part for part in normalized.split("/") if part not in {"", "."}]
                if normalized.startswith("/") or any(part == ".." for part in parts):
                    raise ArtifactError("archive path traversal detected")
                total += max(0, int(info.file_size))
                if total > max_uncompressed:
                    raise ArtifactError("archive expands beyond safety limit")
                compressed = max(1, int(info.compress_size))
                if info.file_size > 1024 * 1024 and info.file_size / compressed > max_ratio:
                    raise ArtifactError("archive compression ratio exceeds safety limit")
    except zipfile.BadZipFile as exc:
        raise ArtifactError("invalid office container") from exc


def validate_blob(data: bytes, fmt: str, max_bytes: int) -> None:
    if not data:
        raise ArtifactError("empty files are not accepted")
    if len(data) > max_bytes:
        raise ArtifactError("file exceeds configured size limit")
    if fmt == "pdf" and not data.startswith(b"%PDF-"):
        raise ArtifactError("file extension does not match PDF content")
    if fmt in OFFICE_FORMATS:
        if not data.startswith(b"PK"):
            raise ArtifactError("file extension does not match Office content")
        _validate_zip_bytes(data)
    if fmt in TEXT_FORMATS:
        try:
            text = data.decode("utf-8-sig")
        except UnicodeDecodeError as exc:
            raise ArtifactError("text file must be UTF-8") from exc
        if fmt == "json":
            try:
                json.loads(text)
            except json.JSONDecodeError as exc:
                raise ArtifactError("invalid JSON document") from exc


def extract_text(path: Path, fmt: str, *, max_chars: int = 120_000) -> tuple[str, dict[str, Any]]:
    metadata: dict[str, Any] = {"format": fmt}
    if fmt in {"txt", "md"}:
        text = path.read_text(encoding="utf-8-sig")
    elif fmt == "json":
        obj = json.loads(path.read_text(encoding="utf-8-sig"))
        text = json.dumps(obj, ensure_ascii=False, indent=2)
        metadata["top_level_type"] = type(obj).__name__
    elif fmt == "csv":
        raw = path.read_text(encoding="utf-8-sig")
        rows = list(csv.reader(io.StringIO(raw)))
        metadata["rows"] = len(rows)
        metadata["columns"] = max((len(row) for row in rows), default=0)
        text = "\n".join(" | ".join(cell for cell in row) for row in rows)
    elif fmt == "pdf":
        if PdfReader is not None:
            reader = PdfReader(str(path))
            metadata["pages"] = len(reader.pages)
            chunks = []
            for index, page in enumerate(reader.pages, start=1):
                chunks.append(f"[Page {index}]\n{page.extract_text() or ''}")
            text = "\n\n".join(chunks)
            if not text.strip():
                metadata["ocr_required"] = True
        else:
            text, extra = _extract_pdf_text(path, max_chars=max_chars)
            metadata.update(extra)
    elif fmt == "docx":
        if Document is not None:
            doc = Document(str(path))
            chunks = [p.text for p in doc.paragraphs if p.text]
            for table in doc.tables:
                for row in table.rows:
                    chunks.append(" | ".join(cell.text for cell in row.cells))
            metadata["paragraphs"] = len(doc.paragraphs)
            metadata["tables"] = len(doc.tables)
            text = "\n".join(chunks)
        else:
            text, extra = _extract_fallback_office_text(path, marker="docx", max_chars=max_chars)
            metadata.update(extra)
    elif fmt == "xlsx":
        text, metadata = _extract_xlsx_text(path, max_chars=max_chars)
    elif fmt == "pptx":
        if Presentation is not None:
            prs = Presentation(str(path))
            metadata["slides"] = len(prs.slides)
            chunks = []
            for index, slide in enumerate(prs.slides, start=1):
                chunks.append(f"[Slide {index}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and str(shape.text).strip():
                        chunks.append(str(shape.text))
            text = "\n".join(chunks)
        else:
            text, extra = _extract_fallback_office_text(path, marker="pptx", max_chars=max_chars)
            metadata.update(extra)
    else:
        raise ArtifactError("unsupported file format")
    if len(text) > max_chars:
        metadata["text_truncated"] = True
        text = text[:max_chars]
    metadata["extracted_chars"] = len(text)
    return text, metadata


def _plain_text_content(content: Any) -> str:
    if isinstance(content, str):
        return content
    if content is None:
        return ""
    return json.dumps(content, ensure_ascii=False, indent=2)


def create_document_bytes(fmt: str, content: Any, *, font_path: str = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf") -> bytes:
    if fmt == "txt":
        return _plain_text_content(content).encode("utf-8")
    if fmt == "md":
        return _plain_text_content(content).encode("utf-8")
    if fmt == "json":
        if isinstance(content, str):
            try:
                content = json.loads(content)
            except json.JSONDecodeError:
                content = {"content": content}
        return json.dumps(content, ensure_ascii=False, indent=2).encode("utf-8")
    if fmt == "csv":
        output = io.StringIO(newline="")
        writer = csv.writer(output)
        if isinstance(content, dict) and isinstance(content.get("rows"), list):
            headers = content.get("headers")
            if isinstance(headers, list):
                writer.writerow(headers)
            for row in content["rows"]:
                writer.writerow(row if isinstance(row, list) else [row])
        elif isinstance(content, list):
            for row in content:
                writer.writerow(row if isinstance(row, list) else [row])
        else:
            writer.writerow([_plain_text_content(content)])
        return output.getvalue().encode("utf-8-sig")
    if fmt == "pdf":
        if SimpleDocTemplate is not None and A4 is not None and mm is not None and pdfmetrics is not None and TTFont is not None and ParagraphStyle is not None and getSampleStyleSheet is not None and Paragraph is not None and Spacer is not None:
            buffer = io.BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=A4, rightMargin=18 * mm, leftMargin=18 * mm, topMargin=18 * mm, bottomMargin=18 * mm)
            font_name = "Helvetica"
            if font_path and Path(font_path).exists():
                font_name = "PARDejaVu"
                if font_name not in pdfmetrics.getRegisteredFontNames():
                    pdfmetrics.registerFont(TTFont(font_name, font_path))
            base = getSampleStyleSheet()["BodyText"]
            style = ParagraphStyle("PARBody", parent=base, fontName=font_name, fontSize=10.5, leading=15)
            story = []
            lines = _plain_text_content(content).splitlines() or [""]
            for line in lines:
                story.append(Paragraph(escape(line) or "&nbsp;", style))
                story.append(Spacer(1, 2.5 * mm))
            doc.build(story)
            return buffer.getvalue()
        return _build_pdf_bytes(content)
    if fmt == "docx":
        if Document is not None:
            doc = Document()
            if isinstance(content, dict):
                title = str(content.get("title") or "").strip()
                if title:
                    doc.add_heading(title, level=1)
                paragraphs = content.get("paragraphs")
                if isinstance(paragraphs, list):
                    for item in paragraphs:
                        doc.add_paragraph(str(item))
                elif "content" in content:
                    for line in str(content.get("content") or "").splitlines():
                        doc.add_paragraph(line)
            else:
                for line in _plain_text_content(content).splitlines() or [""]:
                    doc.add_paragraph(line)
            output = io.BytesIO(); doc.save(output); return output.getvalue()
        return _build_fallback_office_bytes(content, marker="docx")
    if fmt == "xlsx":
        return _build_xlsx_bytes(content)
    if fmt == "pptx":
        if Presentation is not None:
            prs = Presentation()
            title = "Personal Agent Rus"
            paragraphs: list[str]
            if isinstance(content, dict):
                title = str(content.get("title") or title)
                raw = content.get("slides")
                if isinstance(raw, list) and raw:
                    for idx, slide_data in enumerate(raw):
                        layout = prs.slide_layouts[1] if len(prs.slides) or idx else prs.slide_layouts[0]
                        slide = prs.slides.add_slide(layout)
                        if idx == 0 and layout == prs.slide_layouts[0]:
                            slide.shapes.title.text = str(slide_data.get("title") if isinstance(slide_data, dict) else title)
                            if len(slide.placeholders) > 1:
                                slide.placeholders[1].text = str(slide_data.get("content", "") if isinstance(slide_data, dict) else slide_data)
                        else:
                            slide.shapes.title.text = str(slide_data.get("title", f"Slide {idx+1}") if isinstance(slide_data, dict) else f"Slide {idx+1}")
                            if len(slide.placeholders) > 1:
                                slide.placeholders[1].text = str(slide_data.get("content", "") if isinstance(slide_data, dict) else slide_data)
                    output = io.BytesIO(); prs.save(output); return output.getvalue()
                paragraphs = [str(x) for x in content.get("paragraphs", [])] if isinstance(content.get("paragraphs"), list) else [str(content.get("content") or "")]
            else:
                paragraphs = _plain_text_content(content).splitlines()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            slide.placeholders[1].text = "\n".join(paragraphs)
            output = io.BytesIO(); prs.save(output); return output.getvalue()
        return _build_fallback_office_bytes(content, marker="pptx")
    raise ArtifactError("unsupported output format")


class ArtifactService:
    def __init__(self, db_path: Path, root: Path, *, max_bytes: int = 20 * 1024 * 1024):
        self.db_path = Path(db_path)
        self.root = Path(root)
        self.max_bytes = int(max_bytes)
        self.lock = threading.RLock()
        self.root.mkdir(parents=True, exist_ok=True)

    def _db(self) -> Any:
        return connect_app_db(self.db_path)

    def init_schema(self) -> None:
        with self.lock, self._db() as conn:
            conn.executescript(
                """
                CREATE TABLE IF NOT EXISTS artifacts (
                  id TEXT PRIMARY KEY,
                  tenant_id TEXT NOT NULL,
                  user_id TEXT NOT NULL,
                  parent_id TEXT,
                  version INTEGER NOT NULL DEFAULT 1,
                  kind TEXT NOT NULL,
                  original_name TEXT NOT NULL,
                  format TEXT NOT NULL,
                  mime TEXT NOT NULL,
                  storage_key TEXT NOT NULL,
                  size INTEGER NOT NULL,
                  sha256 TEXT NOT NULL,
                  extracted_text TEXT NOT NULL DEFAULT '',
                  metadata_json TEXT NOT NULL DEFAULT '{}',
                  validation_status TEXT NOT NULL,
                  created_at INTEGER NOT NULL,
                  updated_at INTEGER NOT NULL
                );
                CREATE INDEX IF NOT EXISTS idx_artifacts_user_created ON artifacts(user_id, created_at DESC);
                CREATE INDEX IF NOT EXISTS idx_artifacts_parent ON artifacts(parent_id, version);
                """
            )
            conn.commit()

    def _user_root(self, user_id: str) -> Path:
        safe = re.sub(r"[^A-Za-z0-9._-]", "_", str(user_id))[:100] or "unknown"
        root = (self.root / safe / "objects").resolve()
        root.mkdir(parents=True, exist_ok=True)
        if self.root.resolve() not in root.parents:
            raise ArtifactError("invalid workspace path")
        return root

    def _record(self, row: sqlite3.Row | dict[str, Any], *, include_text: bool = False) -> dict[str, Any]:
        item = dict(row)
        result = {
            "artifact_id": item["id"], "tenant_id": item["tenant_id"], "user_id": item["user_id"],
            "parent_id": item.get("parent_id"), "version": item["version"], "kind": item["kind"],
            "name": item["original_name"], "format": item["format"], "mime": item["mime"],
            "size": item["size"], "sha256": item["sha256"], "validation_status": item["validation_status"],
            "created_at": item["created_at"], "updated_at": item["updated_at"],
            "metadata": json.loads(item.get("metadata_json") or "{}"),
            "download_url": f"/api/files/{item['id']}/download",
        }
        if include_text:
            result["text"] = item.get("extracted_text") or ""
        return result

    def _insert_bytes(self, user_id: str, name: str, fmt: str, data: bytes, *, kind: str, parent_id: str | None = None, version: int = 1) -> dict[str, Any]:
        validate_blob(data, fmt, self.max_bytes)
        artifact_id = uuid.uuid4().hex
        display_name = clean_display_name(name, f"artifact.{fmt}")
        if Path(display_name).suffix.lower().lstrip(".") != fmt:
            display_name = f"{Path(display_name).stem or 'artifact'}.{fmt}"
        storage_key = f"{artifact_id}.{fmt}"
        target = self._user_root(user_id) / storage_key
        temp = target.with_suffix(target.suffix + ".tmp")
        temp.write_bytes(data)
        os.replace(temp, target)
        try:
            text, metadata = extract_text(target, fmt)
        except Exception as exc:
            target.unlink(missing_ok=True)
            raise ArtifactError(f"document validation failed: {type(exc).__name__}: {exc}") from exc
        digest = hashlib.sha256(data).hexdigest()
        ts = now_ts()
        mime = MIME_BY_FORMAT[fmt]
        with self.lock, self._db() as conn:
            conn.execute(
                "INSERT INTO artifacts(id,tenant_id,user_id,parent_id,version,kind,original_name,format,mime,storage_key,size,sha256,extracted_text,metadata_json,validation_status,created_at,updated_at) VALUES(?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)",
                (artifact_id, "personal", user_id, parent_id, version, kind, display_name, fmt, mime, storage_key, len(data), digest, text, json.dumps(metadata, ensure_ascii=False), "verified", ts, ts),
            )
            conn.commit()
            row = conn.execute("SELECT * FROM artifacts WHERE id=?", (artifact_id,)).fetchone()
        return self._record(row, include_text=True)

    def upload(self, user_id: str, name: str, data: bytes, requested_format: str = "") -> dict[str, Any]:
        fmt = normalize_format(name, requested_format)
        return self._insert_bytes(user_id, name, fmt, data, kind="upload")

    def create(self, user_id: str, fmt: str, name: str, content: Any) -> dict[str, Any]:
        fmt = normalize_format(name or f"artifact.{fmt}", fmt)
        data = create_document_bytes(fmt, content)
        return self._insert_bytes(user_id, name or f"artifact.{fmt}", fmt, data, kind="generated")

    def get(self, user_id: str, artifact_id: str, *, include_text: bool = True) -> dict[str, Any] | None:
        with self.lock, self._db() as conn:
            row = conn.execute("SELECT * FROM artifacts WHERE id=? AND user_id=?", (artifact_id, user_id)).fetchone()
        return self._record(row, include_text=include_text) if row else None

    def list(self, user_id: str, limit: int = 100) -> list[dict[str, Any]]:
        limit = max(1, min(200, int(limit)))
        with self.lock, self._db() as conn:
            rows = conn.execute("SELECT * FROM artifacts WHERE user_id=? ORDER BY created_at DESC,id DESC LIMIT ?", (user_id, limit)).fetchall()
        return [self._record(row, include_text=False) for row in rows]

    def update(self, user_id: str, artifact_id: str, content: Any, *, name: str = "") -> dict[str, Any]:
        with self.lock, self._db() as conn:
            source = conn.execute("SELECT * FROM artifacts WHERE id=? AND user_id=?", (artifact_id, user_id)).fetchone()
            if not source:
                raise ArtifactError("artifact not found")
            version = int(conn.execute("SELECT COALESCE(MAX(version),0)+1 FROM artifacts WHERE user_id=? AND (id=? OR parent_id=? OR parent_id=?)", (user_id, artifact_id, artifact_id, source["parent_id"] or artifact_id)).fetchone()[0])
        root_parent = source["parent_id"] or source["id"]
        fmt = str(source["format"])
        target_name = name or str(source["original_name"])
        data = create_document_bytes(fmt, content)
        return self._insert_bytes(user_id, target_name, fmt, data, kind="revision", parent_id=root_parent, version=version)

    def delete(self, user_id: str, artifact_id: str) -> bool:
        with self.lock, self._db() as conn:
            row = conn.execute("SELECT storage_key FROM artifacts WHERE id=? AND user_id=?", (artifact_id, user_id)).fetchone()
            if not row:
                return False
            conn.execute("DELETE FROM artifacts WHERE id=? AND user_id=?", (artifact_id, user_id))
            conn.commit()
        (self._user_root(user_id) / row["storage_key"]).unlink(missing_ok=True)
        return True

    def download(self, user_id: str, artifact_id: str) -> tuple[dict[str, Any], Path] | None:
        with self.lock, self._db() as conn:
            row = conn.execute("SELECT * FROM artifacts WHERE id=? AND user_id=?", (artifact_id, user_id)).fetchone()
        if not row:
            return None
        path = (self._user_root(user_id) / row["storage_key"]).resolve()
        if self._user_root(user_id) not in path.parents or not path.is_file():
            return None
        return self._record(row, include_text=False), path

    def contexts(self, user_id: str, artifact_ids: list[str], *, max_total_chars: int = 120_000) -> list[dict[str, Any]]:
        result = []
        used = 0
        for artifact_id in artifact_ids[:12]:
            item = self.get(user_id, artifact_id, include_text=True)
            if not item:
                raise ArtifactError("one or more attached files are unavailable")
            text = str(item.get("text") or "")
            remaining = max_total_chars - used
            if remaining <= 0:
                break
            text = text[:remaining]
            used += len(text)
            result.append({"artifact_id": artifact_id, "name": item["name"], "format": item["format"], "text": text, "sha256": item["sha256"]})
        return result
