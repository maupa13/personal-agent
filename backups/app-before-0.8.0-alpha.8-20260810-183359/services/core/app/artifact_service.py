from __future__ import annotations

import csv
import hashlib
import io
import json
import mimetypes
import os
import re
import sqlite3
import threading
import time
import uuid
import zipfile
from html import escape
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pypdf import PdfReader
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

SUPPORTED_FORMATS = {"txt", "md", "json", "csv", "pdf", "docx", "xlsx", "pptx"}
OFFICE_FORMATS = {"docx", "xlsx", "pptx"}
TEXT_FORMATS = {"txt", "md", "json", "csv"}
MIME_BY_FORMAT = {
    "txt": "text/plain; charset=utf-8",
    "md": "text/markdown; charset=utf-8",
    "json": "application/json",
    "csv": "text/csv; charset=utf-8",
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
EXT_BY_FORMAT = {fmt: f".{fmt}" for fmt in SUPPORTED_FORMATS}
CONTROL_RE = re.compile(r"[\x00-\x1f\x7f]")


class ArtifactError(ValueError):
    pass


def now_ts() -> int:
    return int(time.time())


def clean_display_name(name: str, fallback: str = "file") -> str:
    name = CONTROL_RE.sub("", str(name or "")).strip().replace("\\", "/").split("/")[-1]
    name = name.strip(" .")[:180]
    return name or fallback


def normalize_format(name: str, requested: str = "") -> str:
    requested = str(requested or "").strip().lower().lstrip(".")
    if requested in SUPPORTED_FORMATS:
        return requested
    suffix = Path(clean_display_name(name)).suffix.lower().lstrip(".")
    if suffix in SUPPORTED_FORMATS:
        return suffix
    raise ArtifactError("unsupported file format")


def _validate_zip_bytes(data: bytes, *, max_entries: int = 3000, max_uncompressed: int = 80 * 1024 * 1024, max_ratio: int = 200) -> None:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            infos = zf.infolist()
            if len(infos) > max_entries:
                raise ArtifactError("archive contains too many entries")
            total = 0
            for info in infos:
                normalized = info.filename.replace("\\", "/")
                parts = [part for part in normalized.split("/") if part not in {"", "."}]
                if normalized.startswith("/") or any(part == ".." for part in parts):
                    raise ArtifactError("archive path traversal detected")
                total += max(0, int(info.file_size))
                if total > max_uncompressed:
                    raise ArtifactError("archive expands beyond safety limit")
                compressed = max(1, int(info.compress_size))
                if info.file_size > 1024 * 1024 and info.file_size / compressed > max_ratio:
                    raise ArtifactError("archive compression ratio exceeds safety limit")
    except zipfile.BadZipFile as exc:
        raise ArtifactError("invalid office container") from exc


def validate_blob(data: bytes, fmt: str, max_bytes: int) -> None:
    if not data:
        raise ArtifactError("empty files are not accepted")
    if len(data) > max_bytes:
        raise ArtifactError("file exceeds configured size limit")
    if fmt == "pdf" and not data.startswith(b"%PDF-"):
        raise ArtifactError("file extension does not match PDF content")
    if fmt in OFFICE_FORMATS:
        if not data.startswith(b"PK"):
            raise ArtifactError("file extension does not match Office content")
        _validate_zip_bytes(data)
    if fmt in TEXT_FORMATS:
        try:
            text = data.decode("utf-8-sig")
        except UnicodeDecodeError as exc:
            raise ArtifactError("text file must be UTF-8") from exc
        if fmt == "json":
            try:
                json.loads(text)
            except json.JSONDecodeError as exc:
                raise ArtifactError("invalid JSON document") from exc


def extract_text(path: Path, fmt: str, *, max_chars: int = 120_000) -> tuple[str, dict[str, Any]]:
    metadata: dict[str, Any] = {"format": fmt}
    if fmt in {"txt", "md"}:
        text = path.read_text(encoding="utf-8-sig")
    elif fmt == "json":
        obj = json.loads(path.read_text(encoding="utf-8-sig"))
        text = json.dumps(obj, ensure_ascii=False, indent=2)
        metadata["top_level_type"] = type(obj).__name__
    elif fmt == "csv":
        raw = path.read_text(encoding="utf-8-sig")
        rows = list(csv.reader(io.StringIO(raw)))
        metadata["rows"] = len(rows)
        metadata["columns"] = max((len(row) for row in rows), default=0)
        text = "\n".join(" | ".join(cell for cell in row) for row in rows)
    elif fmt == "pdf":
        reader = PdfReader(str(path))
        metadata["pages"] = len(reader.pages)
        chunks = []
        for index, page in enumerate(reader.pages, start=1):
            chunks.append(f"[Страница {index}]\n{page.extract_text() or ''}")
        text = "\n\n".join(chunks)
        if not text.strip():
            metadata["ocr_required"] = True
    elif fmt == "docx":
        doc = Document(str(path))
        chunks = [p.text for p in doc.paragraphs if p.text]
        for table in doc.tables:
            for row in table.rows:
                chunks.append(" | ".join(cell.text for cell in row.cells))
        metadata["paragraphs"] = len(doc.paragraphs)
        metadata["tables"] = len(doc.tables)
        text = "\n".join(chunks)
    elif fmt == "xlsx":
        wb = load_workbook(str(path), read_only=True, data_only=False)
        metadata["sheets"] = wb.sheetnames
        chunks: list[str] = []
        for ws in wb.worksheets:
            chunks.append(f"[Лист: {ws.title}]")
            row_count = 0
            for row in ws.iter_rows(values_only=False):
                values: list[str] = []
                if any(cell.value is not None for cell in row):
                    row_count += 1
                    for cell in row:
                        value = cell.value
                        if value is None:
                            values.append("")
                        elif getattr(cell, "data_type", "") == "f":
                            values.append(str(value) if str(value).startswith("=") else f"={value}")
                        else:
                            values.append(str(value))
                    chunks.append(" | ".join(values))
                if len("\n".join(chunks)) >= max_chars:
                    break
            metadata.setdefault("sheet_rows", {})[ws.title] = row_count
            if len("\n".join(chunks)) >= max_chars:
                break
        wb.close()
        text = "\n".join(chunks)
    elif fmt == "pptx":
        prs = Presentation(str(path))
        metadata["slides"] = len(prs.slides)
        chunks = []
        for index, slide in enumerate(prs.slides, start=1):
            chunks.append(f"[Слайд {index}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and str(shape.text).strip():
                    chunks.append(str(shape.text))
        text = "\n".join(chunks)
    else:
        raise ArtifactError("unsupported file format")
    if len(text) > max_chars:
        metadata["text_truncated"] = True
        text = text[:max_chars]
    metadata["extracted_chars"] = len(text)
    return text, metadata


def _plain_text_content(content: Any) -> str:
    if isinstance(content, str):
        return content
    if content is None:
        return ""
    return json.dumps(content, ensure_ascii=False, indent=2)


def create_document_bytes(fmt: str, content: Any, *, font_path: str = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf") -> bytes:
    if fmt == "txt":
        return _plain_text_content(content).encode("utf-8")
    if fmt == "md":
        return _plain_text_content(content).encode("utf-8")
    if fmt == "json":
        if isinstance(content, str):
            try:
                content = json.loads(content)
            except json.JSONDecodeError:
                content = {"content": content}
        return json.dumps(content, ensure_ascii=False, indent=2).encode("utf-8")
    if fmt == "csv":
        output = io.StringIO(newline="")
        writer = csv.writer(output)
        if isinstance(content, dict) and isinstance(content.get("rows"), list):
            headers = content.get("headers")
            if isinstance(headers, list):
                writer.writerow(headers)
            for row in content["rows"]:
                writer.writerow(row if isinstance(row, list) else [row])
        elif isinstance(content, list):
            for row in content:
                writer.writerow(row if isinstance(row, list) else [row])
        else:
            writer.writerow([_plain_text_content(content)])
        return output.getvalue().encode("utf-8-sig")
    if fmt == "pdf":
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4, rightMargin=18 * mm, leftMargin=18 * mm, topMargin=18 * mm, bottomMargin=18 * mm)
        font_name = "Helvetica"
        if font_path and Path(font_path).exists():
            font_name = "PARDejaVu"
            if font_name not in pdfmetrics.getRegisteredFontNames():
                pdfmetrics.registerFont(TTFont(font_name, font_path))
        base = getSampleStyleSheet()["BodyText"]
        style = ParagraphStyle("PARBody", parent=base, fontName=font_name, fontSize=10.5, leading=15)
        story = []
        lines = _plain_text_content(content).splitlines() or [""]
        for line in lines:
            story.append(Paragraph(escape(line) or "&nbsp;", style))
            story.append(Spacer(1, 2.5 * mm))
        doc.build(story)
        return buffer.getvalue()
    if fmt == "docx":
        doc = Document()
        if isinstance(content, dict):
            title = str(content.get("title") or "").strip()
            if title:
                doc.add_heading(title, level=1)
            paragraphs = content.get("paragraphs")
            if isinstance(paragraphs, list):
                for item in paragraphs:
                    doc.add_paragraph(str(item))
            elif "content" in content:
                for line in str(content.get("content") or "").splitlines():
                    doc.add_paragraph(line)
        else:
            for line in _plain_text_content(content).splitlines() or [""]:
                doc.add_paragraph(line)
        output = io.BytesIO(); doc.save(output); return output.getvalue()
    if fmt == "xlsx":
        wb = Workbook(); ws = wb.active; ws.title = "Data"
        if isinstance(content, dict) and isinstance(content.get("rows"), list):
            headers = content.get("headers")
            if isinstance(headers, list):
                ws.append(headers)
            for row in content["rows"]:
                ws.append(row if isinstance(row, list) else [row])
        elif isinstance(content, list):
            for row in content:
                ws.append(row if isinstance(row, list) else [row])
        else:
            for line in _plain_text_content(content).splitlines() or [""]:
                ws.append([line])
        output = io.BytesIO(); wb.save(output); wb.close(); return output.getvalue()
    if fmt == "pptx":
        prs = Presentation()
        title = "Personal Agent Rus"
        paragraphs: list[str]
        if isinstance(content, dict):
            title = str(content.get("title") or title)
            raw = content.get("slides")
            if isinstance(raw, list) and raw:
                for idx, slide_data in enumerate(raw):
                    layout = prs.slide_layouts[1] if len(prs.slides) or idx else prs.slide_layouts[0]
                    slide = prs.slides.add_slide(layout)
                    if idx == 0 and layout == prs.slide_layouts[0]:
                        slide.shapes.title.text = str(slide_data.get("title") if isinstance(slide_data, dict) else title)
                        if len(slide.placeholders) > 1:
                            slide.placeholders[1].text = str(slide_data.get("content", "") if isinstance(slide_data, dict) else slide_data)
                    else:
                        slide.shapes.title.text = str(slide_data.get("title", f"Слайд {idx+1}") if isinstance(slide_data, dict) else f"Слайд {idx+1}")
                        if len(slide.placeholders) > 1:
                            slide.placeholders[1].text = str(slide_data.get("content", "") if isinstance(slide_data, dict) else slide_data)
                output = io.BytesIO(); prs.save(output); return output.getvalue()
            paragraphs = [str(x) for x in content.get("paragraphs", [])] if isinstance(content.get("paragraphs"), list) else [str(content.get("content") or "")]
        else:
            paragraphs = _plain_text_content(content).splitlines()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = "\n".join(paragraphs)
        output = io.BytesIO(); prs.save(output); return output.getvalue()
    raise ArtifactError("unsupported output format")


class ArtifactService:
    def __init__(self, db_path: Path, root: Path, *, max_bytes: int = 20 * 1024 * 1024):
        self.db_path = Path(db_path)
        self.root = Path(root)
        self.max_bytes = int(max_bytes)
        self.lock = threading.RLock()
        self.root.mkdir(parents=True, exist_ok=True)

    def _db(self) -> sqlite3.Connection:
        conn = sqlite3.connect(self.db_path, timeout=15, check_same_thread=False)
        conn.row_factory = sqlite3.Row
        return conn

    def init_schema(self) -> None:
        with self.lock, self._db() as conn:
            conn.executescript(
                """
                CREATE TABLE IF NOT EXISTS artifacts (
                  id TEXT PRIMARY KEY,
                  tenant_id TEXT NOT NULL,
                  user_id TEXT NOT NULL,
                  parent_id TEXT,
                  version INTEGER NOT NULL DEFAULT 1,
                  kind TEXT NOT NULL,
                  original_name TEXT NOT NULL,
                  format TEXT NOT NULL,
                  mime TEXT NOT NULL,
                  storage_key TEXT NOT NULL,
                  size INTEGER NOT NULL,
                  sha256 TEXT NOT NULL,
                  extracted_text TEXT NOT NULL DEFAULT '',
                  metadata_json TEXT NOT NULL DEFAULT '{}',
                  validation_status TEXT NOT NULL,
                  created_at INTEGER NOT NULL,
                  updated_at INTEGER NOT NULL
                );
                CREATE INDEX IF NOT EXISTS idx_artifacts_user_created ON artifacts(user_id, created_at DESC);
                CREATE INDEX IF NOT EXISTS idx_artifacts_parent ON artifacts(parent_id, version);
                """
            )
            conn.commit()

    def _user_root(self, user_id: str) -> Path:
        safe = re.sub(r"[^A-Za-z0-9._-]", "_", str(user_id))[:100] or "unknown"
        root = (self.root / safe / "objects").resolve()
        root.mkdir(parents=True, exist_ok=True)
        if self.root.resolve() not in root.parents:
            raise ArtifactError("invalid workspace path")
        return root

    def _record(self, row: sqlite3.Row | dict[str, Any], *, include_text: bool = False) -> dict[str, Any]:
        item = dict(row)
        result = {
            "artifact_id": item["id"], "tenant_id": item["tenant_id"], "user_id": item["user_id"],
            "parent_id": item.get("parent_id"), "version": item["version"], "kind": item["kind"],
            "name": item["original_name"], "format": item["format"], "mime": item["mime"],
            "size": item["size"], "sha256": item["sha256"], "validation_status": item["validation_status"],
            "created_at": item["created_at"], "updated_at": item["updated_at"],
            "metadata": json.loads(item.get("metadata_json") or "{}"),
            "download_url": f"/api/files/{item['id']}/download",
        }
        if include_text:
            result["text"] = item.get("extracted_text") or ""
        return result

    def _insert_bytes(self, user_id: str, name: str, fmt: str, data: bytes, *, kind: str, parent_id: str | None = None, version: int = 1) -> dict[str, Any]:
        validate_blob(data, fmt, self.max_bytes)
        artifact_id = uuid.uuid4().hex
        display_name = clean_display_name(name, f"artifact.{fmt}")
        if Path(display_name).suffix.lower().lstrip(".") != fmt:
            display_name = f"{Path(display_name).stem or 'artifact'}.{fmt}"
        storage_key = f"{artifact_id}.{fmt}"
        target = self._user_root(user_id) / storage_key
        temp = target.with_suffix(target.suffix + ".tmp")
        temp.write_bytes(data)
        os.replace(temp, target)
        try:
            text, metadata = extract_text(target, fmt)
        except Exception as exc:
            target.unlink(missing_ok=True)
            raise ArtifactError(f"document validation failed: {type(exc).__name__}: {exc}") from exc
        digest = hashlib.sha256(data).hexdigest()
        ts = now_ts()
        mime = MIME_BY_FORMAT[fmt]
        with self.lock, self._db() as conn:
            conn.execute(
                "INSERT INTO artifacts(id,tenant_id,user_id,parent_id,version,kind,original_name,format,mime,storage_key,size,sha256,extracted_text,metadata_json,validation_status,created_at,updated_at) VALUES(?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)",
                (artifact_id, "personal", user_id, parent_id, version, kind, display_name, fmt, mime, storage_key, len(data), digest, text, json.dumps(metadata, ensure_ascii=False), "verified", ts, ts),
            )
            conn.commit()
            row = conn.execute("SELECT * FROM artifacts WHERE id=?", (artifact_id,)).fetchone()
        return self._record(row, include_text=True)

    def upload(self, user_id: str, name: str, data: bytes, requested_format: str = "") -> dict[str, Any]:
        fmt = normalize_format(name, requested_format)
        return self._insert_bytes(user_id, name, fmt, data, kind="upload")

    def create(self, user_id: str, fmt: str, name: str, content: Any) -> dict[str, Any]:
        fmt = normalize_format(name or f"artifact.{fmt}", fmt)
        data = create_document_bytes(fmt, content)
        return self._insert_bytes(user_id, name or f"artifact.{fmt}", fmt, data, kind="generated")

    def get(self, user_id: str, artifact_id: str, *, include_text: bool = True) -> dict[str, Any] | None:
        with self.lock, self._db() as conn:
            row = conn.execute("SELECT * FROM artifacts WHERE id=? AND user_id=?", (artifact_id, user_id)).fetchone()
        return self._record(row, include_text=include_text) if row else None

    def list(self, user_id: str, limit: int = 100) -> list[dict[str, Any]]:
        limit = max(1, min(200, int(limit)))
        with self.lock, self._db() as conn:
            rows = conn.execute("SELECT * FROM artifacts WHERE user_id=? ORDER BY created_at DESC,id DESC LIMIT ?", (user_id, limit)).fetchall()
        return [self._record(row, include_text=False) for row in rows]

    def update(self, user_id: str, artifact_id: str, content: Any, *, name: str = "") -> dict[str, Any]:
        with self.lock, self._db() as conn:
            source = conn.execute("SELECT * FROM artifacts WHERE id=? AND user_id=?", (artifact_id, user_id)).fetchone()
            if not source:
                raise ArtifactError("artifact not found")
            version = int(conn.execute("SELECT COALESCE(MAX(version),0)+1 FROM artifacts WHERE user_id=? AND (id=? OR parent_id=? OR parent_id=?)", (user_id, artifact_id, artifact_id, source["parent_id"] or artifact_id)).fetchone()[0])
        root_parent = source["parent_id"] or source["id"]
        fmt = str(source["format"])
        target_name = name or str(source["original_name"])
        data = create_document_bytes(fmt, content)
        return self._insert_bytes(user_id, target_name, fmt, data, kind="revision", parent_id=root_parent, version=version)

    def delete(self, user_id: str, artifact_id: str) -> bool:
        with self.lock, self._db() as conn:
            row = conn.execute("SELECT storage_key FROM artifacts WHERE id=? AND user_id=?", (artifact_id, user_id)).fetchone()
            if not row:
                return False
            conn.execute("DELETE FROM artifacts WHERE id=? AND user_id=?", (artifact_id, user_id))
            conn.commit()
        (self._user_root(user_id) / row["storage_key"]).unlink(missing_ok=True)
        return True

    def download(self, user_id: str, artifact_id: str) -> tuple[dict[str, Any], Path] | None:
        with self.lock, self._db() as conn:
            row = conn.execute("SELECT * FROM artifacts WHERE id=? AND user_id=?", (artifact_id, user_id)).fetchone()
        if not row:
            return None
        path = (self._user_root(user_id) / row["storage_key"]).resolve()
        if self._user_root(user_id) not in path.parents or not path.is_file():
            return None
        return self._record(row, include_text=False), path

    def contexts(self, user_id: str, artifact_ids: list[str], *, max_total_chars: int = 120_000) -> list[dict[str, Any]]:
        result = []
        used = 0
        for artifact_id in artifact_ids[:12]:
            item = self.get(user_id, artifact_id, include_text=True)
            if not item:
                raise ArtifactError("one or more attached files are unavailable")
            text = str(item.get("text") or "")
            remaining = max_total_chars - used
            if remaining <= 0:
                break
            text = text[:remaining]
            used += len(text)
            result.append({"artifact_id": artifact_id, "name": item["name"], "format": item["format"], "text": text, "sha256": item["sha256"]})
        return result
